import os
import re
import uuid
import threading
import subprocess
import tempfile
import yaml
import httpx
import requests
from dotenv import load_dotenv
from ollama import Client as OllamaClient
from datetime import date
from flask import Flask, request, jsonify, send_from_directory, send_file

load_dotenv(os.path.join(os.path.dirname(os.path.abspath(__file__)), '.env'))

# ── ReportLab ──────────────────────────────────────────────────────────────
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.lib.enums import TA_LEFT, TA_CENTER
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

# ── Paths ───────────────────────────────────────────────────────────────────
BASE_DIR     = os.path.dirname(os.path.abspath(__file__))
FRONTEND_DIR = os.path.join(BASE_DIR, '..', 'frontend')
OUTPUT_FILE  = os.path.join(BASE_DIR, 'text.txt')
PDF_DIR      = os.path.join(BASE_DIR, 'pdfs')
PROMPTS_FILE = os.path.join(BASE_DIR, '..', 'prompts', 'cs_claude_sonnet.yaml')

os.makedirs(PDF_DIR, exist_ok=True)

# ── Unicode font (Czech characters) ─────────────────────────────────────────
_FONT_CANDIDATES = [
    # DejaVu — Linux (various distros)
    '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
    '/usr/share/fonts/truetype/ttf-dejavu/DejaVuSans.ttf',
    '/usr/share/fonts/dejavu-sans-fonts/DejaVuSans.ttf',
    # DejaVu — macOS (Homebrew)
    '/opt/homebrew/share/fonts/dejavu-fonts/DejaVuSans.ttf',
    '/usr/local/share/fonts/DejaVuSans.ttf',
    # Arial — macOS system (supports Latin Extended / Czech)
    '/System/Library/Fonts/Supplemental/Arial.ttf',
    '/Library/Fonts/Arial.ttf',
    # Arial Unicode — macOS (full Unicode fallback)
    '/System/Library/Fonts/Supplemental/Arial Unicode.ttf',
    '/Library/Fonts/Arial Unicode.ttf',
]
_BOLD_CANDIDATES = [
    # DejaVu Bold — Linux
    '/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf',
    '/usr/share/fonts/truetype/ttf-dejavu/DejaVuSans-Bold.ttf',
    '/usr/share/fonts/dejavu-sans-fonts/DejaVuSans-Bold.ttf',
    # DejaVu Bold — macOS (Homebrew)
    '/opt/homebrew/share/fonts/dejavu-fonts/DejaVuSans-Bold.ttf',
    '/usr/local/share/fonts/DejaVuSans-Bold.ttf',
    # Arial Bold — macOS system
    '/System/Library/Fonts/Supplemental/Arial Bold.ttf',
    '/Library/Fonts/Arial Bold.ttf',
    # Fallback to regular Arial Unicode (no separate bold file)
    '/System/Library/Fonts/Supplemental/Arial Unicode.ttf',
    '/Library/Fonts/Arial Unicode.ttf',
]

def _first_existing(paths):
    for p in paths:
        if os.path.exists(p):
            return p
    return None

_reg_path  = _first_existing(_FONT_CANDIDATES)
_bold_path = _first_existing(_BOLD_CANDIDATES)

if _reg_path:
    pdfmetrics.registerFont(TTFont('Body',     _reg_path))
    pdfmetrics.registerFont(TTFont('Body-Bold', _bold_path or _reg_path))
    BODY_FONT = 'Body'
    BOLD_FONT = 'Body-Bold'
else:
    BODY_FONT = 'Helvetica'
    BOLD_FONT = 'Helvetica-Bold'

# ── Prompts ──────────────────────────────────────────────────────────────────
with open(PROMPTS_FILE, 'r', encoding='utf-8') as _f:
    PROMPTS = {p['id']: p for p in yaml.safe_load(_f)['prompts']}

QWEN_PROMPTS_FILE = os.path.join(BASE_DIR, '..', 'prompts', 'promptQWEN.yaml')
with open(QWEN_PROMPTS_FILE, 'r', encoding='utf-8') as _f:
    _raw = _f.read()
# Extract the template between the === PROMPT PRO QWEN3:14B === and === KONEC PROMPTU === markers
_m = re.search(r'={3,}\nPROMPT PRO QWEN3:14B\n={3,}\n(.*?)={3,}\nKONEC PROMPTU', _raw, re.DOTALL)
QWEN_TEMPLATE = _m.group(1).strip() if _m else _raw

# ── KKY University Ollama ────────────────────────────────────────────────────
OLLAMA_HOST  = os.environ['OLLAMA_HOST']
OLLAMA_USER  = os.environ['OLLAMA_USER']
OLLAMA_PASS  = os.environ['OLLAMA_PASS']
OLLAMA_MODEL = os.environ.get('OLLAMA_MODEL', 'qwen3:14b')
AI_ENABLED   = True  # always available — university server

# ── KKY ASR ─────────────────────────────────────────────────────────────────
UWEBASR_BASE_URL = 'https://uwebasr.zcu.cz/api/v2/lindat'
LANGUAGE_MODELS = {
    'cs': 'generic/cs/zipformer',
    'sk': 'generic/sk/zipformer',
    'en': 'generic/en/zipformer',
    'de': 'generic/de/zipformer',
    'pl': 'generic/pl/zipformer',
    'hu': 'generic/hu/zipformer',
    'hr': 'generic/hr/zipformer',
    'sr': 'generic/sr/zipformer',
    'nl': 'generic/nl',
}

# ── ffmpeg audio extraction ──────────────────────────────────────────────────

def to_wav(file_data: bytes, filename: str = '') -> bytes:
    """Convert any ffmpeg-supported media to mono 16 kHz WAV bytes."""
    suffix = os.path.splitext(filename)[1].lower() or '.bin'
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as src:
        src.write(file_data)
        src_path = src.name
    dst_path = src_path + '.wav'
    try:
        subprocess.run(
            ['ffmpeg', '-y', '-i', src_path,
             '-vn', '-acodec', 'pcm_s16le', '-ar', '16000', '-ac', '1',
             dst_path],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        with open(dst_path, 'rb') as f:
            return f.read()
    finally:
        os.unlink(src_path)
        if os.path.exists(dst_path):
            os.unlink(dst_path)


# ── In-memory task store ─────────────────────────────────────────────────────
tasks: dict[str, dict] = {}

# ── Flask ────────────────────────────────────────────────────────────────────
app = Flask(__name__, static_folder=FRONTEND_DIR)
app.config['MAX_CONTENT_LENGTH'] = 500 * 1024 * 1024


# ─────────────────────────────────────────────────────────────────────────────
# PDF generation
# ─────────────────────────────────────────────────────────────────────────────

def _md_to_html(text: str) -> str:
    """Convert **bold** markdown to reportlab HTML <b> tags."""
    return re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', text)


def _build_pdf_styles():
    base = getSampleStyleSheet()
    common = dict(fontName=BODY_FONT, leading=16)
    bold   = dict(fontName=BOLD_FONT, leading=16)
    return {
        'title':     ParagraphStyle('title',     fontName=BOLD_FONT, fontSize=18, leading=24, spaceAfter=6,  alignment=TA_CENTER),
        'subtitle':  ParagraphStyle('subtitle',  fontName=BODY_FONT, fontSize=11, leading=14, spaceAfter=20, alignment=TA_CENTER, textColor='#666666'),
        'section':   ParagraphStyle('section',   fontName=BOLD_FONT, fontSize=13, leading=18, spaceBefore=18, spaceAfter=6),
        'body':      ParagraphStyle('body',      fontSize=10, spaceAfter=4,  **common),
        'list_item': ParagraphStyle('list_item', fontSize=10, spaceAfter=3, leftIndent=12, **common),
    }


def _section_to_flowables(text: str, styles: dict) -> list:
    """Parse a Claude markdown response into reportlab Flowables."""
    items = []
    for raw_line in text.split('\n'):
        line = raw_line.strip()
        if not line:
            items.append(Spacer(1, 3))
            continue

        html = _md_to_html(line)

        # Headings produced by Claude ("Shrnutí hlavních bodů" etc.)
        if re.match(r'^#{1,3}\s', line):
            heading = re.sub(r'^#+\s*', '', html)
            items.append(Paragraph(heading, styles['section']))

        # Numbered list  "1. **Term** – text"
        elif re.match(r'^\d+\.', line):
            content = re.sub(r'^\d+\.\s*', '', html)
            items.append(Paragraph(f'<b>{_extract_number(line)}.</b> {content}', styles['list_item']))

        # Bullet list  "- **Term** – text"
        elif line.startswith(('- ', '* ')):
            content = html[2:]
            items.append(Paragraph(f'• {content}', styles['list_item']))

        else:
            items.append(Paragraph(html, styles['body']))

    return items


def _extract_number(line: str) -> str:
    m = re.match(r'^(\d+)', line)
    return m.group(1) if m else ''


def generate_pdf(task_id: str, subject_code: str, sections: dict[str, str]) -> str:
    """Build a PDF and return its file path."""
    path = os.path.join(PDF_DIR, f'{task_id}.pdf')
    styles = _build_pdf_styles()

    doc = SimpleDocTemplate(
        path,
        pagesize=A4,
        leftMargin=2.5*cm, rightMargin=2.5*cm,
        topMargin=2.5*cm,  bottomMargin=2.5*cm,
    )

    story = [
        Paragraph('UTEACH.AI', styles['title']),
        Paragraph(f'{subject_code} &nbsp;·&nbsp; {date.today().strftime("%d. %m. %Y")}', styles['subtitle']),
        HRFlowable(width='100%', thickness=1, color='#e2e8f0', spaceAfter=10),
    ]

    section_labels = {
        'lecture_summary': PROMPTS['lecture_summary']['name'],
        'not_in_slides':   PROMPTS['not_in_slides']['name'],
        'glossary':        PROMPTS['glossary']['name'],
    }

    for key, label in section_labels.items():
        if sections.get(key):
            story.append(Spacer(1, 8))
            story.append(Paragraph(label, styles['section']))
            story.append(HRFlowable(width='100%', thickness=0.5, color='#cbd5e1', spaceAfter=6))
            story.extend(_section_to_flowables(sections[key], styles))

    doc.build(story)
    return path


# ─────────────────────────────────────────────────────────────────────────────
# Claude summarisation
# ─────────────────────────────────────────────────────────────────────────────

def call_claude(prompt_id: str, subject_code: str, transcript: str) -> str:
    """Call the university Ollama instance for one prompt."""
    client = OllamaClient(
        host=OLLAMA_HOST,
        auth=httpx.DigestAuth(OLLAMA_USER, OLLAMA_PASS),
    )
    prompt_cfg = PROMPTS[prompt_id]
    system = prompt_cfg['template'].format(subject_code=subject_code)

    response = client.chat(
        model=OLLAMA_MODEL,
        stream=False,
        options={'num_ctx': 8192},
        messages=[
            {'role': 'system', 'content': system},
            {'role': 'user',   'content': f'Přepis přednášky:\n\n{transcript}'},
        ],
    )
    return response['message']['content']


def call_qwen(transcript: str, subject_code: str = 'KKY',
              student_names: str = '', forbidden_terms: str = '',
              presentation_text: str = '') -> str:
    """Single Ollama call using the Czech promptQWEN template (all three sections)."""
    client = OllamaClient(
        host=OLLAMA_HOST,
        auth=httpx.DigestAuth(OLLAMA_USER, OLLAMA_PASS),
    )
    # Use str.replace instead of .format() — transcript may contain C# code with {}
    prompt = (QWEN_TEMPLATE
              .replace('{subject_code}', subject_code)
              .replace('{presentation_text}', presentation_text or '—')
              .replace('{lecture_transcript}', transcript)
              .replace('{forbidden_terms}', forbidden_terms or '—')
              .replace('{student_names}', student_names or '—'))

    response = client.chat(
        model=OLLAMA_MODEL,
        stream=False,
        options={'num_ctx': 16384},
        messages=[
            {'role': 'user', 'content': prompt},
        ],
    )
    return response.message.content


# ─────────────────────────────────────────────────────────────────────────────
# Background processing
# ─────────────────────────────────────────────────────────────────────────────

def process_task(task_id: str, file_data: bytes, filename: str, lang: str, subject_code: str):
    t = tasks[task_id]

    # ── Convert to WAV via ffmpeg (handles MP4, MKV, AVI, MOV, …) ────────────
    try:
        audio_data = to_wav(file_data, filename)
    except subprocess.CalledProcessError as e:
        t.update(status='error', error=f'ffmpeg conversion failed: {e}')
        return

    # ── Transcription ────────────────────────────────────────────────────────
    app_id  = LANGUAGE_MODELS.get(lang, 'generic/cs/zipformer')
    api_url = f'{UWEBASR_BASE_URL}/{app_id}?format=plaintext'
    try:
        resp = requests.post(
            api_url, data=audio_data,
            headers={'Content-Type': 'application/octet-stream'},
            timeout=3600,
        )
        if resp.status_code == 503:
            t.update(status='error', error='All ASR workers are busy. Please try again.')
            return
        if resp.status_code != 200:
            t.update(status='error', error=f'ASR API returned HTTP {resp.status_code}.')
            return
        transcript = resp.text
    except requests.Timeout:
        t.update(status='error', error='ASR request timed out after 1 hour.')
        return
    except requests.RequestException as e:
        t.update(status='error', error=f'Network error: {e}')
        return

    with open(OUTPUT_FILE, 'w', encoding='utf-8') as fh:
        fh.write(transcript)

    t.update(status='done', result=transcript)


# ─────────────────────────────────────────────────────────────────────────────
# Routes
# ─────────────────────────────────────────────────────────────────────────────

@app.route('/health')
def health():
    return jsonify({
        'status': 'ok',
        'ai_enabled': AI_ENABLED,
        'model': OLLAMA_MODEL,
    }), 200


@app.route('/')
def index():
    return send_from_directory(FRONTEND_DIR, 'index.html')


@app.route('/<path:filename>')
def static_files(filename):
    return send_from_directory(FRONTEND_DIR, filename)


@app.route('/transcribe', methods=['POST'])
def transcribe():
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    f = request.files['file']
    if not f.filename:
        return jsonify({'error': 'No file selected'}), 400

    lang         = request.form.get('language', 'cs')
    subject_code = request.form.get('subject_code', 'KKY').strip() or 'KKY'
    filename     = f.filename or ''
    file_data    = f.read()

    task_id = str(uuid.uuid4())
    tasks[task_id] = {
        'status':    'processing',
        'result':    None,
        'summary':   None,
        'pdf_ready': False,
        'error':     None,
    }

    threading.Thread(
        target=process_task,
        args=(task_id, file_data, filename, lang, subject_code),
        daemon=True,
    ).start()

    return jsonify({'task_id': task_id})


@app.route('/status/<task_id>')
def status(task_id):
    task = tasks.get(task_id)
    if task is None:
        return jsonify({'error': 'Task not found'}), 404
    return jsonify(task)


@app.route('/pdf/<task_id>')
def download_pdf(task_id):
    pdf_path = os.path.join(PDF_DIR, f'{task_id}.pdf')
    if not os.path.exists(pdf_path):
        return jsonify({'error': 'PDF not ready'}), 404
    return send_file(pdf_path, as_attachment=True, download_name='lecture_summary.pdf')


@app.route('/generate-pdf', methods=['POST'])
def generate_pdf_route():
    data = request.get_json(force=True)
    summary_text = (data or {}).get('summary', '').strip()
    subject_code = (data or {}).get('subject_code', 'KKY')
    if not summary_text:
        return jsonify({'error': 'No summary provided'}), 400

    # Split the Qwen all-in-one output into the three named sections
    section_keys = {
        'Shrnutí hlavních bodů': 'lecture_summary',
        'Co nenajdete v prezentaci': 'not_in_slides',
        'Důležité pojmy': 'glossary',
    }
    sections: dict[str, str] = {k: '' for k in section_keys.values()}
    parts = re.split(r'(?m)^#{1,3}\s*(' + '|'.join(re.escape(h) for h in section_keys) + r')\s*$', summary_text)
    i = 1
    while i < len(parts) - 1:
        key = section_keys.get(parts[i].strip())
        if key:
            sections[key] = parts[i + 1].strip()
        i += 2
    if not any(sections.values()):
        sections['lecture_summary'] = summary_text  # fallback: dump everything

    pdf_id = str(uuid.uuid4())
    try:
        path = generate_pdf(pdf_id, subject_code, sections)
        return send_file(path, as_attachment=True, download_name='lecture_summary.pdf')
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/parse-presentation', methods=['POST'])
def parse_presentation():
    if 'file' not in request.files:
        return jsonify({'error': 'No file'}), 400
    f = request.files['file']
    ext = os.path.splitext(f.filename or '')[1].lower()
    try:
        if ext in ('.pptx', '.ppt'):
            import io as _io
            from pptx import Presentation as PRS
            prs = PRS(_io.BytesIO(f.read()))
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                lines.append(t)
            return jsonify({'text': '\n'.join(lines)})
        elif ext == '.pdf':
            import io as _io, pdfplumber
            texts = []
            with pdfplumber.open(_io.BytesIO(f.read())) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        texts.append(t)
            return jsonify({'text': '\n'.join(texts)})
        else:
            return jsonify({'error': f'Unsupported format: {ext}'}), 400
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/summarize', methods=['POST'])
def summarize():
    data = request.get_json(force=True)
    transcript = (data or {}).get('transcript', '').strip()
    if not transcript:
        return jsonify({'error': 'No transcript provided'}), 400
    subject_code      = (data or {}).get('subject_code', 'KKY')
    presentation_text = (data or {}).get('presentation_text', '').strip()
    try:
        summary = call_qwen(transcript, subject_code,
                            presentation_text=presentation_text)
        return jsonify({'summary': summary})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/clean', methods=['POST'])
def clean():
    data = request.get_json(force=True)
    transcript    = (data or {}).get('transcript', '')
    student_names = (data or {}).get('student_names', [])
    custom_prompt = (data or {}).get('custom_prompt', '').strip()

    if not student_names and not custom_prompt:
        return jsonify({'cleaned': transcript})

    names_line = ', '.join(student_names) if student_names else ''
    instructions = []
    if names_line:
        instructions.append(f'Anonymizuj tato jména studentů (nahraď je "[Student]"): {names_line}.')
    if custom_prompt:
        instructions.append(custom_prompt)

    client = OllamaClient(
        host=OLLAMA_HOST,
        auth=httpx.DigestAuth(OLLAMA_USER, OLLAMA_PASS),
    )
    system = ('Jsi asistent pro anonymizaci textů. '
              'Uprav přepis přednášky podle následujících instrukcí a vrať pouze upravený text.\n'
              + ' '.join(instructions))
    try:
        response = client.chat(
            model=OLLAMA_MODEL,
            stream=False,
            options={'num_ctx': 8192},
            messages=[
                {'role': 'system', 'content': system},
                {'role': 'user',   'content': transcript},
            ],
        )
        return jsonify({'cleaned': response['message']['content']})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5001, threaded=True)
